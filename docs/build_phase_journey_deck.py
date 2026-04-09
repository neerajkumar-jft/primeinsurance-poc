"""Build a 7-slide STORY ARC deck for PrimeInsurance.

No phase labels. Pure narrative: problem → journey → now → proof → next.
Business audience. 15-minute talk. Speaker notes on every slide.

Structure:
  1. Cover                    30 sec
  2. Where we started         2 min  (the three failures)
  3. What we built            3 min  (one continuous build story)
  4. Where we are now         3 min  (three failures → three answers)
  5. What a business user sees 2.5 min (the app walkthrough)
  6. Proof it works           2 min  (row-count parity + test results)
  7. What's next + thank you  2 min
  Total: ~15 min
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── palette ─────────────────────────────────────────────────────────────────
NAVY     = RGBColor(0x0B, 0x3D, 0x91)
DARK     = RGBColor(0x1A, 0x1A, 0x1A)
GRAY     = RGBColor(0x55, 0x55, 0x55)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xFA)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
BORDER   = RGBColor(0xC8, 0xD4, 0xE8)
RED      = RGBColor(0xB0, 0x2A, 0x2A)
ORANGE   = RGBColor(0xD8, 0x68, 0x20)
GREEN    = RGBColor(0x1F, 0x7A, 0x3E)
TEAL     = RGBColor(0x0D, 0x6E, 0x8C)
GOLD     = RGBColor(0xC9, 0xA4, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ── helpers ─────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill, line=None, line_w=0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def add_oval(slide, x, y, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Helvetica Neue"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for side in ("left", "right", "top", "bottom"):
        setattr(tf, f"margin_{side}", Emu(0))
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font
    return box


def add_bullets(slide, x, y, w, h, items, *, size=13, color=DARK,
                space_after=8, font="Helvetica Neue"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame; tf.word_wrap = True
    for side in ("left", "right", "top", "bottom"):
        setattr(tf, f"margin_{side}", Emu(0))
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = "•  " + item
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = font
    return box


def slide_header(slide, eyebrow, title):
    add_rect(slide, Emu(0), Emu(0), SW, Inches(0.12), NAVY)
    add_text(slide, Inches(0.5), Inches(0.35), Inches(12), Inches(0.35),
             eyebrow, size=11, bold=True, color=NAVY)
    add_text(slide, Inches(0.5), Inches(0.65), Inches(12), Inches(0.9),
             title, size=28, bold=True, color=DARK)


def slide_footer(slide, page, total):
    add_text(slide, Inches(0.5), Inches(7.05), Inches(9), Inches(0.3),
             "PrimeInsurance Data Intelligence Platform  ·  Jellyfish Technologies",
             size=9, color=GRAY)
    add_text(slide, Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.3),
             f"{page} / {total}", size=9, color=GRAY, align=PP_ALIGN.RIGHT)


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


TOTAL = 7


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover (30 sec)
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)

add_rect(s, Emu(0), Emu(0), SW, SH, NAVY)

add_text(s, Inches(0.8), Inches(1.0), Inches(12), Inches(0.5),
         "PRIMEINSURANCE DATA INTELLIGENCE PLATFORM",
         size=13, bold=True, color=RGBColor(0xC8, 0xD4, 0xE8))

add_text(s, Inches(0.8), Inches(1.6), Inches(12), Inches(1.3),
         "From data to decisions",
         size=52, bold=True, color=WHITE)

add_text(s, Inches(0.8), Inches(3.0), Inches(12), Inches(0.5),
         "A 15-minute walk through the journey",
         size=18, color=RGBColor(0xC8, 0xD4, 0xE8))

# Four-step arc visualization
arc_y = Inches(4.2); arc_h = Inches(1.6)
steps = [
    ("01", "THE PROBLEM",    "Three business failures"),
    ("02", "WHAT WE BUILT",  "One platform, many assistants"),
    ("03", "WHERE WE ARE",   "Three answers, live today"),
    ("04", "WHAT'S NEXT",    "Scaling to operational plays"),
]
sgap = Inches(0.2)
sw2 = (SW - Inches(1.6) - sgap * 3) / 4
for i, (num, lbl, sub) in enumerate(steps):
    x = Inches(0.8) + (sw2 + sgap) * i
    add_rect(s, x, arc_y, sw2, arc_h, RGBColor(0x17, 0x4E, 0xAB))
    add_text(s, x + Inches(0.2), arc_y + Inches(0.15), sw2 - Inches(0.4), Inches(0.3),
             num, size=10, bold=True, color=RGBColor(0xC8, 0xD4, 0xE8))
    add_text(s, x + Inches(0.2), arc_y + Inches(0.4), sw2 - Inches(0.4), Inches(0.4),
             lbl, size=12, bold=True, color=WHITE)
    add_text(s, x + Inches(0.2), arc_y + Inches(0.85), sw2 - Inches(0.4), Inches(0.6),
             sub, size=11, color=RGBColor(0xC8, 0xD4, 0xE8))
    if i < 3:
        ax = x + sw2 + Inches(0.02)
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                  ax, arc_y + arc_h/2 - Inches(0.08),
                                  sgap - Inches(0.04), Inches(0.16))
        arr.fill.solid(); arr.fill.fore_color.rgb = WHITE
        arr.line.fill.background()

add_text(s, Inches(0.8), Inches(6.7), Inches(12), Inches(0.4),
         "Team: Neeraj Kumar (lead) · Abhinav Sarkar · AK Singh · Paras Dhyani    ·    Jellyfish Technologies    ·    2026-04-09",
         size=10, color=RGBColor(0xC8, 0xD4, 0xE8))

set_notes(s,
    "[30 sec] Good morning. Over the next 15 minutes I'm going to walk you "
    "through the PrimeInsurance journey. Four steps: the problem we were "
    "asked to solve, what we built, where it stands today, and what comes "
    "next. Business lens — no deep technical detours. Let's dive in."
)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Where we started (2 min)
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "01  ·  WHERE WE STARTED", "Three business failures, one deadline")

add_text(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.4),
         "Six acquisitions, six databases, three failures the leadership team couldn't answer.",
         size=14, color=GRAY)

problems = [
    (RED, "THE REGULATOR",
     "\"How many customers\ndo you have?\"",
     ["Same person under 6 different IDs",
      "Count inflated by ~12%",
      "No single source of truth",
      "90-day deadline"]),
    (ORANGE, "THE CLAIMS MANAGER",
     "\"Where's the backlog\nthis morning?\"",
     ["18-day processing time",
      "Industry benchmark: 7 days",
      "3 disconnected systems per claim",
      "One undifferentiated queue"]),
    (GREEN, "THE SALES HEAD",
     "\"Which cars do I\nneed to act on today?\"",
     ["Cars sitting 90+ days in one region",
      "Same models selling in 2 weeks elsewhere",
      "No cross-regional visibility",
      "15–20% revenue at risk"]),
]

pw = Inches(4.05); pgap = Inches(0.15); py = Inches(2.25); ph = Inches(4.4)
for i, (color, role, question, bullets) in enumerate(problems):
    x = Inches(0.5) + (pw + pgap) * i
    add_rect(s, x, py, pw, ph, LIGHT_BG)
    add_rect(s, x, py, pw, Inches(0.12), color)
    add_text(s, x + Inches(0.3), py + Inches(0.3), pw - Inches(0.6), Inches(0.35),
             role, size=11, bold=True, color=color)
    add_text(s, x + Inches(0.3), py + Inches(0.7), pw - Inches(0.6), Inches(1.05),
             question, size=17, bold=True, color=DARK)
    add_bullets(s, x + Inches(0.3), py + Inches(2.0), pw - Inches(0.6), ph - Inches(2.1),
                bullets, size=11, space_after=7)

add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.35),
         "Three different people. Three different mornings. Three questions the data couldn't answer.",
         size=12, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

slide_footer(s, 2, TOTAL)

set_notes(s,
    "[2 min] Three people, three mornings, three failures. The compliance "
    "officer couldn't tell the regulator how many customers PrimeInsurance "
    "had — same person duplicated across six acquired systems, count "
    "inflated twelve percent, ninety-day deadline. The claims manager had "
    "an eighteen-day backlog against a seven-day industry benchmark, "
    "because adjusters had to cross-reference three disconnected systems "
    "per claim. And the sales head watched cars sit ninety days in one "
    "region while the same models sold in two weeks in another. Everything "
    "we built traces back to one of these three."
)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — What we built (3 min) — continuous build, no phase labels
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "02  ·  WHAT WE BUILT", "Three capability layers, one continuous platform")

add_text(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.4),
         "Every piece built on the last. Nothing was thrown away along the way.",
         size=14, color=GRAY)

# Three horizontal capability bands
layers = [
    (NAVY, "ONE SOURCE OF TRUTH",
     "Unified, governed, audit-ready",
     "14 raw files from 6 regional systems  →  one clean star schema.  "
     "3,605 raw customer rows deduplicated to 1,604 real people.  "
     "Quality rules catch every bad record; every catch has a resolution path."),
    (TEAL, "AI ASSISTANTS FOR EVERY AUDIENCE",
     "Four personas, four answer shapes",
     "Compliance reads plain-English data quality explanations.  "
     "Fraud investigators see 128 pre-scored cases with investigation briefs.  "
     "Adjusters get 2-second policy lookups.  "
     "Executives read AI-generated narratives from aggregated KPIs."),
    (ORANGE, "DIRECTLY IN USERS' HANDS",
     "One URL. No SQL. No training.",
     "A live database behind the scenes that a business app can actually query.  "
     "A single app screen with the four numbers each persona cares about.  "
     "A plain-English question box that calls an AI over the data — "
     "tested with 10 real questions, zero wrong answers."),
]
ly = Inches(2.15); lh = Inches(1.45); lgap = Inches(0.15)
for i, (color, title, subtitle, body) in enumerate(layers):
    y = ly + (lh + lgap) * i
    add_rect(s, Inches(0.5), y, SW - Inches(1.0), lh, LIGHT_BG)
    add_rect(s, Inches(0.5), y, Inches(0.12), lh, color)

    # Number badge
    add_oval(s, Inches(0.75), y + Inches(0.35), Inches(0.75), Inches(0.75), color)
    add_text(s, Inches(0.75), y + Inches(0.45), Inches(0.75), Inches(0.55),
             str(i + 1), size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Title + subtitle
    add_text(s, Inches(1.7), y + Inches(0.2), Inches(6.5), Inches(0.4),
             title, size=13, bold=True, color=color)
    add_text(s, Inches(1.7), y + Inches(0.5), Inches(6.5), Inches(0.45),
             subtitle, size=16, bold=True, color=DARK)

    # Body — right side, spans remaining width
    add_text(s, Inches(1.7), y + Inches(1.0), SW - Inches(2.2), Inches(0.4),
             body, size=10, color=GRAY)

# Bottom banner — numbers strip
nby = Inches(6.7); nbh = Inches(0.55)
add_text(s, Inches(0.5), nby, SW - Inches(1.0), nbh,
         "14 source files   ·   13,086 records   ·   1,604 real customers   ·   128 fraud flags   ·   6 serving tables   ·   5 Genie tables   ·   10/10 test questions",
         size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

slide_footer(s, 3, TOTAL)

set_notes(s,
    "[3 min] We built this in three capability layers, each on top of the "
    "last. First, one source of truth — 14 raw files from six regional "
    "systems became one clean star schema, with 3,605 raw customer rows "
    "deduplicated down to 1,604 real people, and quality rules catching "
    "every bad record. Second, AI assistants — one per business audience: "
    "plain-English data quality for compliance, 128 pre-scored fraud cases "
    "with briefs, two-second policy lookups for adjusters, and AI-written "
    "narratives for executives. Third, we put it directly in the users' "
    "hands: a live database, a single app screen, a plain-English question "
    "box. Nothing was thrown away between layers. Every piece still runs."
)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Where we are now (3 min)
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "03  ·  WHERE WE ARE NOW", "Three failures, three answers")

add_text(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.4),
         "All three questions the business couldn't answer at the start are now one click away.",
         size=14, color=GRAY)

# Three "answered" cards matching the problems from slide 2
answers = [
    (RED, "✓", "THE REGULATOR", "\"1,604 customers\"",
     "Deduplication methodology visible and auditable. "
     "Regulator query turnaround: seconds, not days. "
     "Compliance team answers directly, without calling engineering."),
    (ORANGE, "✓", "THE CLAIMS MANAGER", "\"Here's the backlog\"",
     "Rejection rate, fraud flags, severity breakdown — all visible on "
     "one screen before morning standup. 276 rejected claims, 128 scored "
     "for fraud investigation with briefs."),
    (GREEN, "✓", "THE SALES HEAD", "\"176 cars need action\"",
     "Every unsold car labeled Fresh, Watch, Stale, or Critical. "
     "176 cars currently in Stale plus Critical. Visible the moment she "
     "opens the app, sorted by days listed."),
]

aw = Inches(4.05); agap = Inches(0.15); ay = Inches(2.25); ah = Inches(4.4)
for i, (color, check, role, quote, body) in enumerate(answers):
    x = Inches(0.5) + (aw + agap) * i
    add_rect(s, x, ay, aw, ah, LIGHT_BG)
    add_rect(s, x, ay, aw, Inches(0.12), color)

    # Check mark tile
    add_oval(s, x + Inches(0.3), ay + Inches(0.3), Inches(0.65), Inches(0.65), color)
    add_text(s, x + Inches(0.3), ay + Inches(0.35), Inches(0.65), Inches(0.55),
             check, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, x + Inches(1.15), ay + Inches(0.35), aw - Inches(1.4), Inches(0.35),
             role, size=10, bold=True, color=color)
    add_text(s, x + Inches(1.15), ay + Inches(0.6), aw - Inches(1.4), Inches(0.5),
             "SOLVED", size=11, bold=True, color=GREEN)

    add_text(s, x + Inches(0.3), ay + Inches(1.2), aw - Inches(0.6), Inches(1.0),
             quote, size=22, bold=True, color=DARK)
    add_text(s, x + Inches(0.3), ay + Inches(2.5), aw - Inches(0.6), ah - Inches(2.6),
             body, size=11, color=GRAY)

add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.35),
         "Same three people. Same three mornings. Different conversation.",
         size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

slide_footer(s, 4, TOTAL)

set_notes(s,
    "[3 min] Now look at the same three people this morning. Compliance "
    "opens the app and sees 1,604 customers, with the deduplication "
    "methodology visible behind it. Her answer to the regulator takes "
    "seconds. The claims manager sees 276 rejected claims and 128 fraud "
    "flags on the same screen before his standup — no waiting for a "
    "weekly report. And the sales head opens the app and sees 176 aging "
    "cars, each labeled Fresh, Watch, Stale, or Critical. She knows "
    "which ones to act on before the discount window closes. Three "
    "failures, three answers. Same three people, different conversation."
)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — What a business user sees (2.5 min)
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "04  ·  THE APP", "One screen. Four numbers. One question box.")

add_text(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.4),
         "This is what business users open every morning. Everything above is behind the scenes.",
         size=14, color=GRAY)

# Full-width app mock
mx = Inches(0.9); my = Inches(2.2); mw = SW - Inches(1.8); mh = Inches(4.4)
add_rect(s, mx, my, mw, mh, WHITE, line=BORDER, line_w=1)

# Browser bar-ish header
add_rect(s, mx, my, mw, Inches(0.45), NAVY)
add_text(s, mx + Inches(0.3), my + Inches(0.1), mw - Inches(0.6), Inches(0.3),
         "🏢  PrimeInsurance Data Intelligence",
         size=15, bold=True, color=WHITE)

# KPI strip
kpy = my + Inches(0.7); kph = Inches(1.55); kpgap = Inches(0.15)
kpw = (mw - Inches(0.6) - kpgap * 3) / 4
mock_kpis = [
    ("1,604",  "Unique customers",  "For the regulator",     NAVY),
    ("1,000",  "Total claims",      "Volume context",        NAVY),
    ("276",    "Rejected claims",   "Backlog signal",        RED),
    ("176",    "Aging inventory",   "Stale + Critical cars", ORANGE),
]
for i, (n, lbl, sub, color) in enumerate(mock_kpis):
    x = mx + Inches(0.3) + (kpw + kpgap) * i
    add_rect(s, x, kpy, kpw, kph, LIGHT_BG)
    add_rect(s, x, kpy, kpw, Inches(0.08), color)
    add_text(s, x, kpy + Inches(0.2), kpw, Inches(0.65),
             n, size=28, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s, x, kpy + Inches(0.9), kpw, Inches(0.3),
             lbl, size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s, x, kpy + Inches(1.17), kpw, Inches(0.3),
             sub, size=9, color=GRAY, align=PP_ALIGN.CENTER)

# Question box
qby = kpy + kph + Inches(0.3)
add_text(s, mx + Inches(0.3), qby, mw - Inches(0.6), Inches(0.35),
         "Ask a question in plain English", size=13, bold=True, color=DARK)
add_rect(s, mx + Inches(0.3), qby + Inches(0.4), mw - Inches(0.6), Inches(0.45),
         LIGHT_BG, line=BORDER, line_w=1)
add_text(s, mx + Inches(0.45), qby + Inches(0.5), mw - Inches(0.9), Inches(0.3),
         "e.g. Which region has the highest claim volume?", size=11, color=GRAY)
add_rect(s, mx + Inches(0.3), qby + Inches(0.95), Inches(1.3), Inches(0.4),
         RGBColor(0xFF, 0x6B, 0x6B))
add_text(s, mx + Inches(0.3), qby + Inches(0.98), Inches(1.3), Inches(0.35),
         "Ask Genie", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Annotation arrows/labels
add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.35),
         "No SQL. No dashboards to learn. No training. One URL — opens, answers, acts.",
         size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

slide_footer(s, 5, TOTAL)

set_notes(s,
    "[2.5 min] This is the app. One URL, one screen. Top: four numbers the "
    "business cares about. 1,604 unique customers — that's compliance's "
    "answer. 1,000 total claims with 276 rejected — that's the claims "
    "manager's backlog signal. 176 aging inventory — that's the sales "
    "head's action list. Below: a plain-English question box. The user "
    "types whatever they want to know, and the AI answers with data. "
    "No SQL, no dashboard training, no navigating through tabs. This is "
    "the first piece of the POC where a business user can act without "
    "calling engineering. That's the whole point."
)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Proof it works (2 min)
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "05  ·  PROOF IT WORKS", "Not a demo — a tested system")

add_text(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.4),
         "Three kinds of evidence: exact data parity, real question testing, and the trick question that proves the AI is actually steered.",
         size=13, color=GRAY)

# Left: data parity — dev matches prod exactly
lx = Inches(0.5); ly = Inches(2.2); lw = Inches(6.3); lh = Inches(4.5)
add_rect(s, lx, ly, lw, lh, LIGHT_BG)
add_rect(s, lx, ly, Inches(0.1), lh, NAVY)
add_text(s, lx + Inches(0.3), ly + Inches(0.25), lw - Inches(0.6), Inches(0.4),
         "DATA PARITY", size=11, bold=True, color=NAVY)
add_text(s, lx + Inches(0.3), ly + Inches(0.55), lw - Inches(0.6), Inches(0.55),
         "Every number matches.", size=18, bold=True, color=DARK)
add_text(s, lx + Inches(0.3), ly + Inches(1.2), lw - Inches(0.6), Inches(0.35),
         "Same data in dev and in production — proof the pipeline is reproducible.",
         size=11, color=GRAY)

# Parity numbers as small matching tiles
pny = ly + Inches(1.75); pnh = Inches(0.55); pngap = Inches(0.1)
parities = [
    ("1,604",  "customers"),
    ("1,000",  "claims"),
    ("128",    "fraud flags"),
    ("999",    "policies"),
    ("2,500",  "cars"),
]
pnw = (lw - Inches(0.6) - pngap * 4) / 5
for i, (n, lbl) in enumerate(parities):
    x = lx + Inches(0.3) + (pnw + pngap) * i
    add_rect(s, x, pny, pnw, pnh, WHITE, line=BORDER, line_w=1)
    add_text(s, x, pny + Inches(0.05), pnw, Inches(0.3),
             n, size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x, pny + Inches(0.32), pnw, Inches(0.2),
             lbl, size=8, color=GRAY, align=PP_ALIGN.CENTER)

add_text(s, lx + Inches(0.3), ly + Inches(2.55), lw - Inches(0.6), Inches(0.35),
         "dev ↔ prod, table by table, row for row.", size=11, bold=True, color=GREEN)

add_text(s, lx + Inches(0.3), ly + Inches(3.05), lw - Inches(0.6), Inches(0.3),
         "Only difference: one duplicate sales_id caught by the production layer's stricter constraints — a real data quality bug Gold was hiding.",
         size=10, color=GRAY)

# Right: natural-language test results
rx = Inches(6.9); ry = Inches(2.2); rw = SW - Inches(0.5) - rx; rh = Inches(4.5)
add_rect(s, rx, ry, rw, rh, WHITE, line=BORDER, line_w=1)
add_text(s, rx + Inches(0.3), ry + Inches(0.25), rw - Inches(0.6), Inches(0.4),
         "NATURAL LANGUAGE TEST", size=11, bold=True, color=NAVY)
add_text(s, rx + Inches(0.3), ry + Inches(0.55), rw - Inches(0.6), Inches(0.55),
         "10 real questions. Zero wrong.", size=18, bold=True, color=DARK)

# Score tiles
scy = ry + Inches(1.2); sch = Inches(1.1); scgap = Inches(0.1)
scs = [
    ("10/10", "ANSWERED",   GREEN),
    ("7/10",  "DIRECT",     GREEN),
    ("3/10",  "CLARIFIED",  ORANGE),
    ("0/10",  "WRONG",      GREEN),
]
scw = (rw - Inches(0.6) - scgap * 3) / 4
for i, (n, lbl, color) in enumerate(scs):
    x = rx + Inches(0.3) + (scw + scgap) * i
    add_rect(s, x, scy, scw, sch, LIGHT_BG)
    add_rect(s, x, scy, scw, Inches(0.07), color)
    add_text(s, x, scy + Inches(0.18), scw, Inches(0.5),
             n, size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s, x, scy + Inches(0.72), scw, Inches(0.3),
             lbl, size=9, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

# Highlight
add_text(s, rx + Inches(0.3), ry + Inches(2.5), rw - Inches(0.6), Inches(0.3),
         "Highlight: the trick question", size=11, bold=True, color=DARK)
add_text(s, rx + Inches(0.3), ry + Inches(2.85), rw - Inches(0.6), Inches(1.5),
         "We deliberately asked \"how long do claims take to process?\" — a question our source data can't answer because the dates are corrupted. "
         "The AI correctly refused: \"Claim processing dates are corrupted at the source... the platform cannot compute real processing duration.\"\n\n"
         "That's the strongest proof the AI is actually steered by our domain rules, not just generating plausible-sounding answers.",
         size=10, color=GRAY)

slide_footer(s, 6, TOTAL)

set_notes(s,
    "[2 min] Three kinds of evidence. First, data parity: every number "
    "matches exactly between our development environment and production. "
    "1,604 customers, 1,000 claims, 128 fraud flags — identical row "
    "counts. The pipeline is reproducible. Second, we tested the "
    "natural-language layer with ten real business questions. Ten out of "
    "ten answered, zero wrong, seven direct, three polite clarifications. "
    "And third — the proof I'm proudest of — the trick question. We asked "
    "about claim processing time, which our source data can't answer "
    "because the dates are corrupted. The AI correctly refused, quoting "
    "our domain rules. That's the strongest proof it's actually steered, "
    "not hallucinating."
)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — What's next + thank you (2 min)
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "06  ·  WHAT'S NEXT", "Seven more operational plays — same foundation")

add_text(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.4),
         "Nothing we built has to be thrown away to ship the next use case. That's the point of the foundation.",
         size=14, color=GRAY)

# Seven plays in a grid
plays = [
    ("1", "Claims Triage Workbench", "Live case queue with AI investigation briefs — directly attacks the 18-day backlog."),
    ("2", "Customer Golden Record",  "Live match-or-create service — every new record goes through the dedup layer before it lands."),
    ("3", "Regulator Genie Space",   "A separate self-serve interface for compliance with row-level security by jurisdiction."),
    ("4", "Inventory Marketplace",   "Dealer managers approve and reserve cross-regional transfers — ends the 90-day aging problem."),
    ("5", "Real-Time Fraud Cases",   "Flagged claims become tracked cases with status, notes, and outcomes — closing the loop from detection to resolution."),
    ("6", "Renewal Intelligence",    "Proactive retention queue from the risk-segmentation model we already have."),
    ("7", "Agent Copilot",           "The policy Q&A assistant embedded in the call-center workflow with conversation memory."),
]
py2 = Inches(2.25); pgap = Inches(0.15)
pw2 = (SW - Inches(1.0) - pgap * 3) / 4
ph2 = Inches(1.95)
for i, (num, title, desc) in enumerate(plays[:4]):
    x = Inches(0.5) + (pw2 + pgap) * i
    y = py2
    add_rect(s, x, y, pw2, ph2, LIGHT_BG)
    add_rect(s, x, y, Inches(0.08), ph2, NAVY)
    add_oval(s, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.5), NAVY)
    add_text(s, x + Inches(0.2), y + Inches(0.23), Inches(0.5), Inches(0.45),
             num, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.85), y + Inches(0.2), pw2 - Inches(1.05), Inches(0.8),
             title, size=12, bold=True, color=DARK)
    add_text(s, x + Inches(0.2), y + Inches(0.9), pw2 - Inches(0.4), ph2 - Inches(1.0),
             desc, size=10, color=GRAY)

# Second row — 3 items
py3 = py2 + ph2 + Inches(0.15)
for i, (num, title, desc) in enumerate(plays[4:]):
    x = Inches(0.5) + (pw2 + pgap) * i
    y = py3
    add_rect(s, x, y, pw2, ph2, LIGHT_BG)
    add_rect(s, x, y, Inches(0.08), ph2, NAVY)
    add_oval(s, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.5), NAVY)
    add_text(s, x + Inches(0.2), y + Inches(0.23), Inches(0.5), Inches(0.45),
             num, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.85), y + Inches(0.2), pw2 - Inches(1.05), Inches(0.8),
             title, size=12, bold=True, color=DARK)
    add_text(s, x + Inches(0.2), y + Inches(0.9), pw2 - Inches(0.4), ph2 - Inches(1.0),
             desc, size=10, color=GRAY)

# Thank you in the 4th column of second row
thx_x = Inches(0.5) + (pw2 + pgap) * 3
add_rect(s, thx_x, py3, pw2, ph2, NAVY)
add_text(s, thx_x, py3 + Inches(0.4), pw2, Inches(0.6),
         "Thank you.", size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, thx_x, py3 + Inches(1.05), pw2, Inches(0.4),
         "Questions?", size=14, color=RGBColor(0xC8, 0xD4, 0xE8), align=PP_ALIGN.CENTER)
add_text(s, thx_x, py3 + Inches(1.5), pw2, Inches(0.35),
         "Jellyfish Technologies", size=9, color=RGBColor(0xC8, 0xD4, 0xE8), align=PP_ALIGN.CENTER)

slide_footer(s, 7, TOTAL)

set_notes(s,
    "[2 min] What's next. Seven more operational plays sit on the same "
    "foundation, ready to ship in sequence. Claims Triage Workbench goes "
    "after the eighteen-day backlog directly — a live case queue with AI "
    "briefs. Customer Golden Record turns our dedup into a live service "
    "so the regulator's question stays answered forever. A regulator-facing "
    "Genie space with jurisdictional row-level security. An inventory "
    "marketplace for cross-regional transfers. Real-time fraud case "
    "management, renewal intelligence, and a call-center copilot. "
    "Nothing on this list requires new infrastructure. Every piece builds "
    "on what we showed you today. Thank you — I'm happy to take questions."
)


# ── save ─────────────────────────────────────────────────────────────────────
out = "/Users/neerajkumar/Documents/workspace/primeinsurance-poc/docs/PrimeInsurance_Phase_Journey.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}  ·  Structure: story arc, no phase labels  ·  Timing: ~15 min")
